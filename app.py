import io
import re
from pathlib import Path
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR

st.set_page_config(page_title="TXT → PPT", page_icon="🖼️", layout="centered")

st.title("TXT → PPTX generator")
st.caption("Sorok/bekezdések → diák. Egynyelvű és kétnyelvű mód, tetszőleges tipográfia és színek.")

# ---------- Helpers ----------
def read_text_multi_enc(data: bytes) -> str:
    for enc in ("utf-8", "cp1250", "iso-8859-2", "latin2"):
        try:
            return data.decode(enc)
        except Exception:
            continue
    return data.decode("utf-8", errors="replace")

def parse_text(content: str, mode: str, preserve_blanks: bool):
    if mode == "srt":
        blocks = re.split(r"\n\s*\n", content.strip(), flags=re.MULTILINE)
        out = []
        ts_re = re.compile(r"^\s*\d{2}:\d{2}:\d{2},\d{3}\s*-->\s*\d{2}:\d{2}:\d{2},\d{3}\s*$")
        for b in blocks:
            lines = []
            for ln in b.splitlines():
                s = ln.strip()
                if not s or s.isdigit() or ts_re.match(s):
                    continue
                lines.append(s)
            if lines:
                out.append(" ".join(lines))
        return out
    if mode == "para":
        return [b.strip() for b in re.split(r"\n\s*\n", content) if b.strip()]
    if preserve_blanks:
        return content.splitlines()
    return [l.strip() for l in content.splitlines() if l.strip()]

def parse_bilingual(content: str, *, use_blank_as_separator: bool = True, blank_line_as_slide: bool = False):
    """
    Kétnyelvű párosítás:
      - Nem üres sorokat 2-esével párosítjuk: (line1, line2) -> egy dia.
      - Ha blank_line_as_slide=True: az üres sorok **önálló üres diát** jelentenek.
      - Ha use_blank_as_separator=True: az üres sorokat elválasztóként **eldobjuk** (nem zavarják a párosítást).
    Visszatér: list, amelynek elemei:
      - ("", "") -> üres dia jelzés
      - (line1, line2) vagy (line1, "") ha páratlan maradt
    """
    raw_lines = content.splitlines()

    if blank_line_as_slide:
        pairs = []
        buf = []
        for ln in raw_lines:
            s = ln.strip()
            if s == "":
                # üres dia
                pairs.append(("", ""))
                continue
            buf.append(s)
            if len(buf) == 2:
                pairs.append((buf[0], buf[1]))
                buf = []
        if buf:  # páratlan maradt
            pairs.append((buf[0], ""))
        return pairs
    else:
        # nincs üres dia, üreseket dobjuk-e?
        lines = [ln.strip() for ln in raw_lines if (ln.strip() != "" or not use_blank_as_separator)]
        pairs = []
        i = 0
        while i < len(lines):
            l1 = lines[i]
            l2 = lines[i+1] if i+1 < len(lines) else ""
            pairs.append((l1, l2))
            i += 2
        return pairs

def add_background_slide(prs: Presentation, bg_rgb=(0,0,0)):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*bg_rgb)
    return slide

def add_textbox(
    slide, text, *,
    area_left, area_top, area_width, area_height,
    font_name="Arial", font_size_pt=44, font_rgb=(255,255,255),
    bold=False, italic=False,
    align_center=True, shrink_to_fit=True,
    vertical_position="bottom"   # "bottom" vagy "top"
):
    tb = slide.shapes.add_textbox(area_left, area_top, area_width, area_height)
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0

    if vertical_position == "top":
        tf.vertical_anchor = MSO_ANCHOR.TOP
    else:
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM

    if shrink_to_fit:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    p = tf.paragraphs[0]
    p.text = text
    p.font.name = (font_name or "Arial").strip()
    p.font.size = Pt(float(font_size_pt))
    p.font.color.rgb = RGBColor(*font_rgb)
    p.font.bold = bool(bold)
    p.font.italic = bool(italic)
    p.alignment = PP_ALIGN.CENTER if align_center else PP_ALIGN.LEFT

def add_text_slide_single(
    prs: Presentation, text: str, *,
    m_top_cm=0.13, m_bottom_cm=0.13, m_left_cm=0.25, m_right_cm=0.25,
    font_name="Arial", font_size_pt=44, font_rgb=(255,255,255),
    bold=False, italic=False,
    bg_rgb=(0,0,0), align_center=True, shrink_to_fit=True,
    vertical_position="bottom"
):
    slide = add_background_slide(prs, bg_rgb=bg_rgb)
    sw, sh = prs.slide_width, prs.slide_height
    left = Cm(m_left_cm)
    width = sw - (Cm(m_left_cm) + Cm(m_right_cm))
    height = sh - (Cm(m_top_cm) + Cm(m_bottom_cm))
    top = Cm(m_top_cm)

    add_textbox(
        slide, text,
        area_left=left, area_top=top, area_width=width, area_height=height,
        font_name=font_name, font_size_pt=font_size_pt, font_rgb=font_rgb,
        bold=bold, italic=italic,
        align_center=align_center, shrink_to_fit=shrink_to_fit,
        vertical_position=vertical_position
    )

def add_text_slide_bilingual(
    prs: Presentation, line1: str, line2: str, *,
    m_left_cm=0.25, m_right_cm=0.25,
    text_band_height_cm=2.5,
    primary_offset_cm=0.0,
    secondary_offset_cm=1.6,
    primary_font=("Arial", 44, (255,255,255), False, False),
    secondary_font=("Arial", 36, (200,200,200), False, False),
    bg_rgb=(0,0,0),
    align_center=True,
    shrink_to_fit=True,
    vertical_position="bottom"
):
    slide = add_background_slide(prs, bg_rgb=bg_rgb)
    sw, sh = prs.slide_width, prs.slide_height

    left = Cm(m_left_cm)
    width = sw - (Cm(m_left_cm) + Cm(m_right_cm))
    band_h = Cm(text_band_height_cm)

    if vertical_position == "top":
        # top-aligned: offset a dia tetejétől értendő
        p_top = Cm(primary_offset_cm)
        s_top = Cm(secondary_offset_cm)
    else:
        # bottom-aligned: offset a dia aljától értendő
        p_top = sh - Cm(primary_offset_cm) - band_h
        s_top = sh - Cm(secondary_offset_cm) - band_h

    add_textbox(
        slide, line1,
        area_left=left, area_top=p_top, area_width=width, area_height=band_h,
        font_name=primary_font[0], font_size_pt=primary_font[1], font_rgb=primary_font[2],
        bold=primary_font[3], italic=primary_font[4],
        align_center=align_center, shrink_to_fit=shrink_to_fit,
        vertical_position=vertical_position
    )

    add_textbox(
        slide, line2,
        area_left=left, area_top=s_top, area_width=width, area_height=band_h,
        font_name=secondary_font[0], font_size_pt=secondary_font[1], font_rgb=secondary_font[2],
        bold=secondary_font[3], italic=secondary_font[4],
        align_center=align_center, shrink_to_fit=shrink_to_fit,
        vertical_position=vertical_position
    )

def build_ppt(
    items,
    *,
    widescreen=True,
    mode="single",
    shrink_to_fit=True,
    blank_slide_on_empty=False,
    # single mode
    m_top_cm=0.13, m_bottom_cm=0.13, m_left_cm=0.25, m_right_cm=0.25,
    font_name="Arial", font_size_pt=44, font_rgb=(255,255,255),
    single_bold=False, single_italic=False,
    # bilingual
    bottom_band_height_cm=2.5,
    primary_bottom_offset_cm=0.0, secondary_bottom_offset_cm=1.6,
    primary_font=("Arial", 44, (255,255,255), False, False),
    secondary_font=("Arial", 36, (200,200,200), False, False),
    # common
    bg_rgb=(0,0,0),
    align_center=True,
    vertical_position="bottom"
) -> bytes:
    prs = Presentation()
    if widescreen:
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

    if mode == "bilingual":
        for item in items:
            if isinstance(item, tuple):
                l1, l2 = item
                add_text_slide_bilingual(
                    prs, l1, l2,
                    m_left_cm=m_left_cm, m_right_cm=m_right_cm,
                    text_band_height_cm=bottom_band_height_cm,
                    primary_offset_cm=primary_bottom_offset_cm,
                    secondary_offset_cm=secondary_bottom_offset_cm,
                    primary_font=primary_font, secondary_font=secondary_font,
                    bg_rgb=bg_rgb, align_center=align_center, shrink_to_fit=shrink_to_fit,
                    vertical_position=vertical_position
                )
            else:
                # item == ("", "") eset helyett egy stringes jelzés is lehet – itt üres dia
                add_background_slide(prs, bg_rgb=bg_rgb)
    else:
        for it in items:
            if blank_slide_on_empty and (it.strip() == ""):
                add_background_slide(prs, bg_rgb=bg_rgb)
            else:
                add_text_slide_single(
                    prs, it,
                    m_top_cm=m_top_cm, m_bottom_cm=m_bottom_cm, m_left_cm=m_left_cm, m_right_cm=m_right_cm,
                    font_name=font_name, font_size_pt=font_size_pt, font_rgb=font_rgb,
                    bold=single_bold, italic=single_italic,
                    bg_rgb=bg_rgb, align_center=align_center, shrink_to_fit=shrink_to_fit,
                    vertical_position=vertical_position
                )

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio

# ---------- UI ----------
st.subheader("Forrás")
tab1, tab2 = st.tabs(["📄 Fájl feltöltés", "📂 Meglévő útvonal"])

uploaded = None
with tab1:
    f = st.file_uploader("Válassz .txt vagy .srt fájlt", type=["txt", "srt"])
    if f is not None:
        data = f.read()
        text = read_text_multi_enc(data)
        uploaded = ("uploaded", text)

with tab2:
    p = st.text_input("Teljes elérési út")
    if p:
        try:
            data = Path(p).read_bytes()
            text = read_text_multi_enc(data)
            uploaded = ("path", text)
        except Exception as e:
            st.error(f"Nem tudtam beolvasni: {e}")

st.subheader("Beállítások")

mode_col, fmt_col = st.columns([1,2])

with mode_col:
    bilingual = st.checkbox("Bilingual mode (2 sor/dia)", value=False)
    st.write("")
    
    caption_position = st.radio(
    "Felirat pozíciója",
    ["bottom", "top"],
    index=0,
    format_func=lambda x: "Dia alján" if x == "bottom" else "Dia tetején"
)
    widescreen = st.checkbox("Widescreen 16:9", value=True)
    shrink = st.checkbox("Hosszú sorok tördelése", value=True)
    align_center = st.checkbox("Középre igazítás", value=True)
    bg_hex = st.color_picker("Háttér szín", "#000000")

with fmt_col:
    if not bilingual:
        st.markdown("#### Egynyelvű tipográfia")
        font_name = st.text_input("Betűtípus", value="Arial")
        font_size_pt = st.number_input("Betűméret (pt)", 8.0, 200.0, 44.0, 1.0)
        font_hex = st.color_picker("Szöveg szín", "#FFFFFF")
        single_bold = st.checkbox("Félkövér", value=False)
        single_italic = st.checkbox("Dőlt", value=False)

        st.markdown("#### Külső margók (cm)")
        m_top_cm = st.number_input("Felső margó", 0.0, 10.0, 1.0, 0.01)
        m_bottom_cm = st.number_input("Alsó margó", 0.0, 10.0, 1.0, 0.01)
        m_left_cm = st.number_input("Bal margó", 0.0, 20.0, 3.0, 0.01)
        m_right_cm = st.number_input("Jobb margó", 0.0, 20.0, 3.0, 0.01)

        st.write("")
        blank_on_empty = st.checkbox("Üres sor → üres dia", value=True)
        st.write("")

    else:
        st.markdown("#### Kétnyelvű tipográfia")
        st.write("***Primer (alsóbb) sor:***")
        prim_font = st.text_input("Primer betűtípus", value="Arial", key="prim_font")
        prim_size = st.number_input("Primer betűméret (pt)", 8.0, 200.0, 44.0, 1.0, key="prim_size")
        prim_hex = st.color_picker("Primer szín", "#FFFFFF", key="prim_hex")
        prim_bold = st.checkbox("Primer félkövér", value=False, key="prim_bold")
        prim_italic = st.checkbox("Primer dőlt", value=True, key="prim_italic")
        prim_offset = st.number_input("Primer alsó offset (cm)", 0.0, 10.0, 1.0, 0.1, key="prim_off")
        st.write("")

        st.write("***Szekunder (fölötte lévő) sor:***")
        sec_font = st.text_input("Szekunder betűtípus", value="Arial", key="sec_font")
        sec_size = st.number_input("Szekunder betűméret (pt)", 8.0, 200.0, 44.0, 1.0, key="sec_size")
        sec_hex = st.color_picker("Szekunder szín", "#C8C8C8", key="sec_hex")
        sec_bold = st.checkbox("Szekunder félkövér", value=False, key="sec_bold")
        sec_italic = st.checkbox("Szekunder dőlt", value=False, key="sec_italic")
        sec_offset = st.number_input("Szekunder alsó offset (cm)", 0.0, 10.0, 5.0, 0.1, key="sec_off")
        st.write("")

        st.markdown("#### Elrendezés")
        bottom_band = st.number_input("Szövegdoboz magasság (cm)", 1.0, 10.0, 4.0, 0.1, key="band_h")
        m_left_cm = st.number_input("Bal margó (cm)", 0.0, 20.0, 3.0, 0.01, key="biml")
        m_right_cm = st.number_input("Jobb margó (cm)", 0.0, 20.0, 3.0, 0.01, key="bimr")

        st.write("")
        use_blank_sep = st.checkbox("Üres sor csak elválasztó", value=False)
        bi_blankline_slide = st.checkbox("Üres sor → üres dia", value=True)
        st.write("")

if st.button("PPTX generálása", type="primary", use_container_width=True):
    if not uploaded:
        st.warning("Adj meg forrást (fájl vagy útvonal).")
    else:
        _, raw_text = uploaded
        bg_rgb = tuple(int(bg_hex.lstrip("#")[i:i+2], 16) for i in (0,2,4))

        if bilingual:
            pairs = parse_bilingual(
                raw_text,
                use_blank_as_separator=use_blank_sep,
                blank_line_as_slide=bi_blankline_slide
            )
            primary_font = (prim_font, prim_size, tuple(int(prim_hex.lstrip("#")[i:i+2], 16) for i in (0,2,4)),
                            prim_bold, prim_italic)
            secondary_font = (sec_font, sec_size, tuple(int(sec_hex.lstrip("#")[i:i+2], 16) for i in (0,2,4)),
                              sec_bold, sec_italic)
            pptx_bytes = build_ppt(
                pairs,
                widescreen=widescreen, mode="bilingual", shrink_to_fit=shrink,
                m_left_cm=m_left_cm, m_right_cm=m_right_cm,
                bottom_band_height_cm=bottom_band,
                primary_bottom_offset_cm=prim_offset,
                secondary_bottom_offset_cm=sec_offset,
                primary_font=primary_font, secondary_font=secondary_font,
                bg_rgb=bg_rgb, align_center=align_center,
                vertical_position=caption_position
            )
            st.success(f"Siker! {len([p for p in pairs if isinstance(p, tuple)])} kétnyelvű dia + "
                       f"{len([p for p in pairs if not isinstance(p, tuple)])} üres dia.")
        else:
            font_rgb = tuple(int(font_hex.lstrip("#")[i:i+2], 16) for i in (0,2,4))
            items = parse_text(raw_text, mode="line", preserve_blanks=blank_on_empty)
            pptx_bytes = build_ppt(
                items,
                widescreen=widescreen, mode="single", shrink_to_fit=shrink,
                blank_slide_on_empty=blank_on_empty,
                m_top_cm=m_top_cm, m_bottom_cm=m_bottom_cm,
                m_left_cm=m_left_cm, m_right_cm=m_right_cm,
                font_name=font_name, font_size_pt=font_size_pt, font_rgb=font_rgb,
                single_bold=single_bold, single_italic=single_italic,
                bg_rgb=bg_rgb, align_center=align_center,
                vertical_position=caption_position
            )
            st.success(f"Siker! {len(items)} egynyelvű dia (az üres sorok külön diát kaphattak).")

        st.download_button(
            "PPTX letöltése",
            data=pptx_bytes,
            file_name="slides.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
